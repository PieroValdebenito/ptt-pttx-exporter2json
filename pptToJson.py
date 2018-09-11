# -*- coding: utf-8 -*-
from pptx import Presentation
import glob

text_runs = []
page = 0
for hymn in glob.glob("hymnsES/*.pptx"):
    prs = Presentation(hymn)
    text = '<h3>'
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text = text +" "+ run.text.replace("strofa","strofa <br>").replace("Coro","<b>Coro</b><br> ")
                if count == 0:
                    text = text + " </h3>"
                count=count+1
                text = text + " " 
            text = text + " <br><br>" 
    text = text + "<br>" 
    text_runs.append(text)
output = "this.lyrics = {"
for t in text_runs:
   page = page + 1
   output = output+'hymn'+str(page)+': "'+ t +'",\n'
output = output.replace("  "," ")
output = output.replace("  "," ")
output = output.replace("<br><br>","<br>")

output = output + "};"
print output