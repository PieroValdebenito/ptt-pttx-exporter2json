# -*- coding: utf-8 -*-
from pptx import Presentation
import glob

text_runs = []
page = 0
for hymn in glob.glob("hymnsES/*.pptx"):
    prs = Presentation(hymn)
    text = '.'
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text = text +" "+ run.text
                    text = text.replace(". Himno   ","")
                    text = text.replace(". Himno  ","")
                    text = text.replace(". Himno ","")  
                    text = text.replace(".  Himno ","")                  

                break
        break
    text_runs.append(text)
output = "this.titles = ["
for t in text_runs:
   page = page + 1
   output = output+"{"+'number: '+str(page)+',  title: "'+ t +'"},\n'
   output = output.replace('" ','"').replace("  "," ")
   output = output.replace('" ','"').replace("  "," ")
output = output + "];"
print output